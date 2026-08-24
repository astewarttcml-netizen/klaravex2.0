#!/usr/bin/env python3
"""Klaravex investor pitch deck — public claims only, no internals."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BG = RGBColor(0x0A, 0x0A, 0x10)
PANEL = RGBColor(0x12, 0x12, 0x1B)
PURPLE = RGBColor(0x8B, 0x7C, 0xFF)
WHITE = RGBColor(0xF5, 0xF5, 0xF8)
MUTED = RGBColor(0xA0, 0xA0, 0xB0)
DIM = RGBColor(0x60, 0x60, 0x70)
GREEN = RGBColor(0x4A, 0xDE, 0x96)

W = Inches(13.333)
H = Inches(7.5)
ASSETS = Path("/home/anthony/Klaravex2.0/docs/deck-assets")
OUT = Path("/home/anthony/Klaravex2.0/docs/Klaravex-Investor-Deck-2026-08.pptx")

TOTAL = 11


def _run(p, text, *, size=18, bold=False, color=WHITE, font="Calibri"):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return r


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    _fill(sh, color)
    return sh


def txt(slide, l, t, w, h, text, *, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    _run(p, text, size=size, bold=bold, color=color)
    return box


def lines(slide, l, t, w, h, parts, *, gap=10):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        _run(p, text, size=size, bold=bold, color=color)
    return box


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, BG)
    return s


def photo(slide, file, l, t, w, h):
    slide.shapes.add_picture(str(ASSETS / file), l, t, w, h)


def kicker(slide, n, label):
    txt(slide, Inches(0.55), Inches(0.42), Inches(10), Inches(0.3), label.upper(), size=13, bold=True, color=PURPLE)
    txt(slide, Inches(12.2), Inches(0.42), Inches(0.6), Inches(0.3), f"{n:02d}", size=13, bold=True, color=DIM, align=PP_ALIGN.RIGHT)


def stat_card(slide, x, y, w, h, big, label, *, big_color=WHITE):
    rect(slide, x, y, w, h, PANEL)
    txt(slide, x + Inches(0.25), y + Inches(0.22), w - Inches(0.5), Inches(0.8), big, size=40, bold=True, color=big_color)
    txt(slide, x + Inches(0.25), y + Inches(1.05), w - Inches(0.5), h - Inches(1.2), label, size=14, color=MUTED)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # ── 1 · Cover ────────────────────────────────────────────────
    s = blank(prs)
    photo(s, "klaravex-com.png", Inches(5.0), 0, Inches(8.35), H)
    rect(s, 0, 0, Inches(5.55), H, BG)
    rect(s, Inches(0.55), Inches(2.02), Inches(0.55), Inches(0.09), PURPLE)
    txt(s, Inches(0.55), Inches(2.3), Inches(4.7), Inches(1.1), "Klaravex", size=56, bold=True, color=WHITE)
    txt(s, Inches(0.55), Inches(3.5), Inches(4.6), Inches(1.5),
        "AI-native managed IT and security for American small business.", size=21, color=MUTED)
    txt(s, Inches(0.55), Inches(6.35), Inches(4.6), Inches(0.35), "Investor briefing · August 2026", size=13, color=DIM)
    txt(s, Inches(0.55), Inches(6.72), Inches(4.6), Inches(0.35), "klaravex.com", size=13, color=PURPLE)

    # ── 2 · Problem ──────────────────────────────────────────────
    s = blank(prs)
    kicker(s, 2, "The problem")
    txt(s, Inches(0.55), Inches(0.9), Inches(12.2), Inches(1.6),
        "Small firms face enterprise-grade threats\nwith no IT department.", size=34, bold=True)
    cols = [
        ("No help at 2 a.m.", "Traditional MSPs run human ticket queues. SMBs wait hours for a password reset and days for real answers."),
        ("Priced for enterprises", "Full-service IT and security is built and priced for 500-seat companies — a 12-person law firm can't buy it."),
        ("Conflicted advice", "Most providers earn vendor commissions, so the 'recommendation' is whatever pays the reseller best."),
    ]
    for i, (h2, b) in enumerate(cols):
        x = Inches(0.55 + i * 4.18)
        rect(s, x, Inches(2.9), Inches(3.95), Inches(3.6), PANEL)
        rect(s, x, Inches(2.9), Inches(3.95), Inches(0.07), PURPLE)
        txt(s, x + Inches(0.28), Inches(3.2), Inches(3.4), Inches(0.85), h2, size=21, bold=True)
        txt(s, x + Inches(0.28), Inches(4.15), Inches(3.4), Inches(2.2), b, size=15, color=MUTED)
    txt(s, Inches(0.55), Inches(6.8), Inches(12.2), Inches(0.4),
        "Insurers, clients, and regulators now demand security readiness from firms this size — the pressure is structural, not optional.",
        size=14, color=DIM)

    # ── 3 · Solution ─────────────────────────────────────────────
    s = blank(prs)
    kicker(s, 3, "The solution")
    txt(s, Inches(0.55), Inches(0.9), Inches(12.2), Inches(1.6),
        "Klara AI answers instantly.\nA senior engineer owns the hard 11%.", size=34, bold=True)
    stat_card(s, Inches(0.55), Inches(2.95), Inches(2.95), Inches(2.3), "89%", "of IT issues resolved by AI — instantly, any hour, any day", big_color=PURPLE)
    stat_card(s, Inches(3.7), Inches(2.95), Inches(2.95), Inches(2.3), "24/7", "AI coverage across every US time zone — no queue, no hold music")
    stat_card(s, Inches(6.85), Inches(2.95), Inches(2.95), Inches(2.3), "2hr", "senior engineer SLA for the cases that need human judgment")
    stat_card(s, Inches(10.0), Inches(2.95), Inches(2.78), Inches(2.3), "$0", "vendor commissions — advice with no reseller conflict", big_color=GREEN)
    txt(s, Inches(0.55), Inches(5.6), Inches(12.2), Inches(1.2),
        "The AI is always labeled — clients know when they're talking to Klara AI and when a certified engineer has taken over. "
        "That transparency is the trust wedge against both chatbot-washing and legacy MSPs.",
        size=16, color=MUTED)

    # ── 4 · Product: business ────────────────────────────────────
    s = blank(prs)
    kicker(s, 4, "Product · Business")
    txt(s, Inches(0.55), Inches(0.85), Inches(7.5), Inches(0.55), "Full-stack IT department, per-seat price", size=28, bold=True)
    txt(s, Inches(9.4), Inches(0.95), Inches(3.4), Inches(0.35), "klaravex.com", size=14, color=PURPLE, align=PP_ALIGN.RIGHT)
    photo(s, "klaravex-proof.png", Inches(0.55), Inches(1.6), Inches(8.0), Inches(5.0))
    feats = [
        "Four pillars, fourteen services — network & security, strategy & vCIO, cloud & productivity, infrastructure & support",
        "Readiness for HIPAA, SOC 2, and cyber-insurance — what insurers now require of small firms",
        "Client portal: live system status, AI chat, named engineer, auto-documentation",
    ]
    for i, f in enumerate(feats):
        y = Inches(1.75 + i * 1.65)
        rect(s, Inches(8.85), y, Inches(0.09), Inches(1.35), PURPLE)
        txt(s, Inches(9.1), y, Inches(3.7), Inches(1.5), f, size=14, color=MUTED)

    # ── 5 · Product: consumer ────────────────────────────────────
    s = blank(prs)
    kicker(s, 5, "Product · Consumer")
    txt(s, Inches(0.55), Inches(0.85), Inches(7.5), Inches(0.55), "The same AI, sold to households", size=28, bold=True)
    txt(s, Inches(8.9), Inches(0.95), Inches(3.9), Inches(0.35), "personal.klaravex.com", size=14, color=PURPLE, align=PP_ALIGN.RIGHT)
    photo(s, "personal-klaravex.png", Inches(0.55), Inches(1.6), Inches(8.0), Inches(5.0))
    feats = [
        "Plain-English tech help from $29 a session — no truck roll, so AI keeps the price a fraction of a house call",
        "Family and senior plans, scam recovery offered free — a trust engine that seeds the brand",
        "Consumer track doubles as top-of-funnel: the household buyer is often the small-business owner",
    ]
    for i, f in enumerate(feats):
        y = Inches(1.75 + i * 1.65)
        rect(s, Inches(8.85), y, Inches(0.09), Inches(1.35), PURPLE)
        txt(s, Inches(9.1), y, Inches(3.7), Inches(1.5), f, size=14, color=MUTED)

    # ── 6 · Business model ───────────────────────────────────────
    s = blank(prs)
    kicker(s, 6, "Business model")
    txt(s, Inches(0.55), Inches(0.9), Inches(12.2), Inches(0.6), "Recurring, per-seat, AI-margin economics", size=32, bold=True)
    tiers = [
        ("Foundation", "$49", "per user / month", "Lean teams that need IT to just work"),
        ("Assurance", "$79", "per user / month", "Firms with data, payments, or insurance exposure"),
        ("Directive", "$129", "per user / month", "Regulated firms that need audit-ready security"),
    ]
    for i, (name, price, unit, who) in enumerate(tiers):
        x = Inches(0.55 + i * 3.02)
        rect(s, x, Inches(1.85), Inches(2.85), Inches(3.4), PANEL)
        if name == "Directive":
            rect(s, x, Inches(1.85), Inches(2.85), Inches(0.07), PURPLE)
        txt(s, x + Inches(0.25), Inches(2.1), Inches(2.4), Inches(0.4), name, size=18, bold=True, color=PURPLE)
        txt(s, x + Inches(0.25), Inches(2.55), Inches(2.4), Inches(0.8), price, size=40, bold=True)
        txt(s, x + Inches(0.25), Inches(3.4), Inches(2.4), Inches(0.35), unit, size=13, color=DIM)
        txt(s, x + Inches(0.25), Inches(3.85), Inches(2.4), Inches(1.3), who, size=14, color=MUTED)
    rect(s, Inches(9.75), Inches(1.85), Inches(3.03), Inches(3.4), PANEL)
    txt(s, Inches(10.0), Inches(2.1), Inches(2.6), Inches(0.4), "Consumer", size=18, bold=True, color=GREEN)
    txt(s, Inches(10.0), Inches(2.55), Inches(2.6), Inches(2.6),
        "Sessions from $29\nPlans $29–$39 / month\n\nCash-pay, zero CAC when it converts from free scam-recovery help.",
        size=14, color=MUTED)
    txt(s, Inches(0.55), Inches(5.6), Inches(12.2), Inches(1.2),
        "Every tier is delivered AI-first. The marginal cost of the 89% is compute — human cost concentrates only "
        "on the 11% and on sales. That's the structural margin advantage over labor-priced MSPs.",
        size=16, color=MUTED)

    # ── 7 · Market ───────────────────────────────────────────────
    s = blank(prs)
    kicker(s, 7, "Market")
    txt(s, Inches(0.55), Inches(0.9), Inches(12.2), Inches(0.6), "A $71B US market, majority SMB, growing 11% a year", size=32, bold=True)
    stat_card(s, Inches(0.55), Inches(1.85), Inches(3.9), Inches(2.6), "$71B", "US managed services market, 2026 — projected $120B by 2031", big_color=PURPLE)
    stat_card(s, Inches(4.75), Inches(1.85), Inches(3.9), Inches(2.6), "~58%", "of that market is small and medium enterprises — the fastest-growing segment")
    stat_card(s, Inches(8.95), Inches(1.85), Inches(3.83), Inches(2.6), "11%", "CAGR 2026–2031, driven by cyber-insurance mandates and the security skills gap", big_color=GREEN)
    txt(s, Inches(0.55), Inches(4.85), Inches(12.2), Inches(1.5),
        "The demand driver is structural: insurers require 24/7 monitoring for renewal, and SMBs cannot hire the "
        "$165K/yr security engineers this work demands. They must buy it as a service — and AI-native delivery is "
        "the only way to serve a 10-seat firm profitably.",
        size=16, color=MUTED)
    txt(s, Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.3),
        "Source: Mordor Intelligence, United States Managed Services Market (2026).", size=11, color=DIM)

    # ── 8 · Moat ─────────────────────────────────────────────────
    s = blank(prs)
    kicker(s, 8, "Why we win")
    txt(s, Inches(0.55), Inches(0.9), Inches(12.2), Inches(1.05),
        "The company itself runs on the product.", size=34, bold=True)
    txt(s, Inches(0.55), Inches(2.0), Inches(12.2), Inches(0.9),
        "Klaravex OS — our own operating console — runs delivery, growth, and operations with a fleet of AI agents. "
        "The cost structure competitors would need years to copy is our default.",
        size=17, color=MUTED)
    photo(s, "os-product.png", Inches(0.55), Inches(3.0), Inches(12.23), Inches(2.9))
    txt(s, Inches(0.55), Inches(6.15), Inches(12.2), Inches(0.9),
        "Same playbook the clients buy: AI does the volume, humans do the judgment. It compounds — every client "
        "interaction trains better runbooks, every runbook lowers the marginal cost of the next client.",
        size=15, color=MUTED)

    # ── 9 · Go-to-market ─────────────────────────────────────────
    s = blank(prs)
    kicker(s, 9, "Go-to-market")
    txt(s, Inches(0.55), Inches(0.9), Inches(12.2), Inches(0.6), "Vertical wedge, automated demand", size=32, bold=True)
    cols = [
        ("Regulated verticals first", "Law, accounting, and medical practices — where readiness pressure is highest and a named senior engineer closes the deal."),
        ("Always-on demand engine", "AI-researched outbound, daily content, and search — the same automation that runs delivery runs acquisition, at software cost."),
        ("Two-track funnel", "Consumer help builds trust at $29; the same buyer signs their 12-person firm at $79–$129 a seat."),
    ]
    for i, (h2, b) in enumerate(cols):
        x = Inches(0.55 + i * 4.18)
        rect(s, x, Inches(1.9), Inches(3.95), Inches(3.5), PANEL)
        rect(s, x, Inches(1.9), Inches(3.95), Inches(0.07), PURPLE)
        txt(s, x + Inches(0.28), Inches(2.2), Inches(3.4), Inches(0.95), h2, size=20, bold=True)
        txt(s, x + Inches(0.28), Inches(3.25), Inches(3.4), Inches(2.0), b, size=15, color=MUTED)
    txt(s, Inches(0.55), Inches(5.75), Inches(12.2), Inches(0.9),
        "Free offers do the selling: a 45-minute IT assessment for businesses, free scam-recovery help for consumers. "
        "Both produce a documented, named-engineer first impression.",
        size=16, color=MUTED)

    # ── 10 · Where we are ────────────────────────────────────────
    s = blank(prs)
    kicker(s, 10, "Where we are")
    txt(s, Inches(0.55), Inches(0.9), Inches(12.2), Inches(0.6), "Built, live, and selling", size=32, bold=True)
    rows = [
        ("Product live", "Klara AI resolution, senior-engineer escalation, and the client portal are in production — not a prototype."),
        ("Both surfaces shipping", "klaravex.com (business) and personal.klaravex.com (consumer) are live with published pricing."),
        ("Acquisition running daily", "The automated outbound, content, and ads engine is on and producing pipeline every business day."),
        ("Entity & IP", "Klaravex LLC (Wyoming), nationwide US remote delivery, trademark filed (USPTO serial 99856526)."),
    ]
    for i, (h2, b) in enumerate(rows):
        y = Inches(1.85 + i * 1.18)
        rect(s, Inches(0.55), y, Inches(12.23), Inches(1.02), PANEL)
        txt(s, Inches(0.85), y + Inches(0.14), Inches(3.3), Inches(0.75), h2, size=17, bold=True, color=PURPLE)
        txt(s, Inches(4.4), y + Inches(0.14), Inches(8.2), Inches(0.75), b, size=15, color=MUTED)
    txt(s, Inches(0.55), Inches(6.75), Inches(12.2), Inches(0.4),
        "Current traction metrics (clients, MRR, pipeline) shared in the data room / live conversation.",
        size=13, color=DIM)

    # ── 11 · Ask ─────────────────────────────────────────────────
    s = blank(prs)
    rect(s, Inches(0.55), Inches(1.95), Inches(0.55), Inches(0.09), PURPLE)
    txt(s, Inches(0.55), Inches(2.2), Inches(12.2), Inches(1.0),
        "The MSP industry is a labor business.\nWe rebuilt it as a software business.", size=34, bold=True)
    txt(s, Inches(0.55), Inches(4.3), Inches(12.2), Inches(0.9),
        "Raising to scale sales and onboarding across the three launch verticals.  [amount + terms — see one-pager]",
        size=18, color=MUTED)
    txt(s, Inches(0.55), Inches(5.9), Inches(12.2), Inches(0.4), "Anthony Stewart · Founder", size=16, bold=True)
    txt(s, Inches(0.55), Inches(6.3), Inches(12.2), Inches(0.4), "hello@klaravex.com   ·   klaravex.com", size=15, color=PURPLE)

    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    print("Created:", build())
