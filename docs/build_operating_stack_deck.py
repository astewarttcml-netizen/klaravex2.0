#!/usr/bin/env python3
"""Klaravex product pitch — public surfaces only. No internals."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# Public site brand (not the internal OS phosphor theme)
BG = RGBColor(0x08, 0x08, 0x0C)
PANEL = RGBColor(0x10, 0x10, 0x16)
PURPLE = RGBColor(0x8B, 0x7C, 0xFF)
WHITE = RGBColor(0xF4, 0xF4, 0xF7)
MUTED = RGBColor(0x9A, 0x9A, 0xA8)
DIM = RGBColor(0x5C, 0x5C, 0x68)

W = Inches(13.333)
H = Inches(7.5)
ASSETS = Path("/home/anthony/Klaravex2.0/docs/deck-assets")
OUT = Path("/home/anthony/Klaravex2.0/docs/Klaravex-Operating-Stack-2026-08-24.pptx")


def _run(p, text, *, size=18, bold=False, color=WHITE, font="Calibri"):
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return run


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    _fill(sh, color)
    return sh


def _txt(slide, l, t, w, h, text, *, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(anchor, "t"))
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    _run(p, text, size=size, bold=bold, color=color)
    return box


def _lines(slide, l, t, w, h, parts):
    """parts: list of (text, size, bold, color)."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, part in enumerate(parts):
        text, size, bold, color = part
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        _run(p, text, size=size, bold=bold, color=color)
    return box


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, W, H, BG)
    return s


def photo(slide, file, l, t, w, h):
    slide.shapes.add_picture(str(ASSETS / file), l, t, w, h)


def footer(slide, n, total=8):
    _txt(slide, Inches(0.5), Inches(7.15), Inches(8), Inches(0.25), "Klaravex", size=11, color=DIM)
    _txt(slide, Inches(11.4), Inches(7.15), Inches(1.4), Inches(0.25), f"{n}  /  {total}", size=11, color=DIM, align=PP_ALIGN.RIGHT)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1 Cover — product photo
    s = blank(prs)
    photo(s, "klaravex-com.png", Inches(4.2), Inches(0), Inches(9.2), Inches(7.5))
    _rect(s, 0, 0, Inches(5.4), H, BG)
    _txt(s, Inches(0.55), Inches(2.0), Inches(4.5), Inches(0.35), "PRODUCT DECK", size=12, bold=True, color=PURPLE)
    _txt(s, Inches(0.55), Inches(2.45), Inches(4.6), Inches(1.8), "Klaravex", size=54, bold=True, color=WHITE)
    _txt(
        s,
        Inches(0.55),
        Inches(4.35),
        Inches(4.5),
        Inches(1.4),
        "Managed IT, powered by Klara AI.\nTwo sites. One operating console.",
        size=18,
        color=MUTED,
    )
    _txt(s, Inches(0.55), Inches(6.5), Inches(4.5), Inches(0.4), "klaravex.com", size=14, color=PURPLE)

    # 2 What it is
    s = blank(prs)
    footer(s, 2)
    _txt(s, Inches(0.55), Inches(0.45), Inches(12), Inches(0.3), "THE COMPANY", size=12, bold=True, color=PURPLE)
    _txt(s, Inches(0.55), Inches(0.85), Inches(12), Inches(1.1), "AI takes the first call.\nA senior engineer takes the hard ones.", size=32, bold=True, color=WHITE)
    photo(s, "klaravex-proof.png", Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.6))

    # 3 Business site
    s = blank(prs)
    footer(s, 3)
    _txt(s, Inches(0.55), Inches(0.35), Inches(8), Inches(0.28), "KLARAVEX", size=12, bold=True, color=PURPLE)
    _txt(s, Inches(0.55), Inches(0.68), Inches(12), Inches(0.5), "The business site", size=28, bold=True, color=WHITE)
    _txt(s, Inches(8.6), Inches(0.78), Inches(4.2), Inches(0.35), "klaravex.com", size=14, color=MUTED, align=PP_ALIGN.RIGHT)
    photo(s, "klaravex-com.png", Inches(0.5), Inches(1.35), Inches(12.3), Inches(5.55))

    # 4 Services
    s = blank(prs)
    footer(s, 4)
    _txt(s, Inches(0.55), Inches(0.35), Inches(8), Inches(0.28), "KLARAVEX", size=12, bold=True, color=PURPLE)
    _txt(s, Inches(0.55), Inches(0.68), Inches(12), Inches(0.5), "What we sell is on the site", size=28, bold=True, color=WHITE)
    photo(s, "klaravex-services.png", Inches(0.5), Inches(1.35), Inches(12.3), Inches(5.55))

    # 5 Personal
    s = blank(prs)
    footer(s, 5)
    _txt(s, Inches(0.55), Inches(0.35), Inches(8), Inches(0.28), "PERSONAL", size=12, bold=True, color=PURPLE)
    _txt(s, Inches(0.55), Inches(0.68), Inches(12), Inches(0.5), "Home and home-office support", size=28, bold=True, color=WHITE)
    _txt(s, Inches(8.2), Inches(0.78), Inches(4.6), Inches(0.35), "personal.klaravex.com", size=14, color=MUTED, align=PP_ALIGN.RIGHT)
    photo(s, "personal-klaravex.png", Inches(0.5), Inches(1.35), Inches(12.3), Inches(5.55))

    # 6 Klara AI
    s = blank(prs)
    footer(s, 6)
    photo(s, "klaravex-com.png", Inches(5.1), 0, Inches(8.3), H)
    _rect(s, 0, 0, Inches(5.6), H, BG)
    _txt(s, Inches(0.55), Inches(1.7), Inches(4.7), Inches(0.3), "KLARA AI", size=12, bold=True, color=PURPLE)
    _txt(s, Inches(0.55), Inches(2.15), Inches(4.8), Inches(1.6), "The AI on the site.\nAlways labeled.", size=32, bold=True, color=WHITE)
    _txt(
        s,
        Inches(0.55),
        Inches(4.1),
        Inches(4.7),
        Inches(2.0),
        "Klara AI handles everyday IT. A named senior engineer steps in when judgment is required. Clients always know which one they’re talking to.",
        size=16,
        color=MUTED,
    )

    # 7 Klaravex OS
    s = blank(prs)
    footer(s, 7)
    _txt(s, Inches(0.55), Inches(0.4), Inches(12), Inches(0.28), "KLARAVEX OS", size=12, bold=True, color=PURPLE)
    _txt(s, Inches(0.55), Inches(0.75), Inches(12), Inches(0.7), "The operating console", size=32, bold=True, color=WHITE)
    _txt(
        s,
        Inches(0.55),
        Inches(1.5),
        Inches(12.2),
        Inches(0.55),
        "One screen for the company: systems, agents, and communications.",
        size=18,
        color=MUTED,
    )
    photo(s, "os-product.png", Inches(0.5), Inches(2.2), Inches(12.3), Inches(3.95))

    # 8 Klaravex 2.0 + close
    s = blank(prs)
    _txt(s, Inches(0.55), Inches(0.4), Inches(12), Inches(0.28), "KLARAVEX 2.0", size=12, bold=True, color=PURPLE)
    _txt(s, Inches(0.55), Inches(0.8), Inches(12.2), Inches(1.3), "Keeps the sites and the outreach moving.", size=32, bold=True, color=WHITE)
    _txt(
        s,
        Inches(0.55),
        Inches(2.15),
        Inches(12.2),
        Inches(0.7),
        "Content, campaigns, and conversations that point back to the two websites. Not a second product — the growth layer behind this one.",
        size=18,
        color=MUTED,
    )
    photo(s, "klaravex-com.png", Inches(0.5), Inches(3.1), Inches(6.05), Inches(3.55))
    photo(s, "personal-klaravex.png", Inches(6.75), Inches(3.1), Inches(6.05), Inches(3.55))
    footer(s, 8)

    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    print("Created:", build())
