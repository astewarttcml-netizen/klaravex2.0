#!/usr/bin/env python3
"""Klaravex investor deck v2 — flagship rebuild per DECK-SPEC.md.

Design system: full-bleed near-black, indigo->violet gradient accents,
brand logo, big display numerals, shape-built diagrams, one idea per slide.
External voice: no internal codenames, no vendor names, no first person.
"""
from pptx import Presentation
from pptx.util import Inches as IN, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- brand tokens -----------------------------------------------------------
INK      = RGBColor(0x14, 0x14, 0x17)   # page background
INK2     = RGBColor(0x1C, 0x1C, 0x22)   # card on dark
INDIGO   = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_L = RGBColor(0x81, 0x8C, 0xF8)
VIOLET   = RGBColor(0x7C, 0x3A, 0xED)
CREAM    = RGBColor(0xF5, 0xF3, 0xEE)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTE     = RGBColor(0x9A, 0x9A, 0xA6)
GREEN    = RGBColor(0x34, 0xD3, 0x99)

FONT = "Noto Sans"
W, H = IN(13.333), IN(7.5)
LOGO_LIGHT = "/home/anthony/klaravex/brand/exports/klaravex-logo-light-2x.png"
ICON = "/home/anthony/klaravex/brand/exports/klaravex-icon-transparent-400.png"
OUT = "/home/anthony/klaravex/decks-2026-08-24/v2/00-investor-overview-v2.pptx"

prs = Presentation()
prs.slide_width, prs.slide_height = W, H


def slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = INK; r.line.fill.background()
    r.shadow.inherit = False
    return s


def grad_bar(s, x, y, w, h):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.line.fill.background(); sh.shadow.inherit = False
    sh.fill.gradient()
    stops = sh.fill.gradient_stops
    stops[0].color.rgb = INDIGO
    stops[1].color.rgb = VIOLET
    try:
        sh.fill.gradient_angle = 0
    except Exception:
        pass
    return sh


def card(s, x, y, w, h, fill=INK2, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.06
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.shadow.inherit = False
    if line:
        sh.line.color.rgb = line; sh.line.width = Pt(1.2)
    else:
        sh.line.fill.background()
    return sh


def txt(s, x, y, w, h, parts, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_before=6):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (t, size, bold, color) in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if i:
            p.space_before = Pt(space_before)
        r = p.add_run(); r.text = t
        f = r.font
        f.name = FONT; f.size = Pt(size); f.bold = bold; f.color.rgb = color
    return tb


def kicker(s, text_):
    grad_bar(s, IN(0.7), IN(0.62), IN(0.5), IN(0.075))
    txt(s, IN(0.7), IN(0.78), IN(11.9), IN(0.5),
        [(text_.upper(), 13, True, INDIGO_L)])


def headline(s, text_, size=33):
    txt(s, IN(0.7), IN(1.18), IN(12.0), IN(1.2), [(text_, size, True, WHITE)])


def pagefoot(s, n):
    txt(s, IN(12.35), IN(7.02), IN(0.75), IN(0.35),
        [(f"{n:02d}", 11, False, MUTE)], align=PP_ALIGN.RIGHT)
    s.shapes.add_picture(ICON, IN(0.7), IN(6.95), height=IN(0.32))


def arrow(s, x, y, w=IN(0.42)):
    sh = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, IN(0.3))
    sh.fill.solid(); sh.fill.fore_color.rgb = INDIGO
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


# ---- 1 · cover --------------------------------------------------------------
s = slide()
grad_bar(s, 0, H - IN(0.14), W, IN(0.14))
s.shapes.add_picture(LOGO_LIGHT, IN(0.7), IN(0.85), height=IN(0.55))
txt(s, IN(0.7), IN(2.75), IN(11.9), IN(2.2),
    [("Managed IT, run by AI.", 54, True, WHITE),
     ("Priced like it.", 54, True, INDIGO_L)], space_before=0)
txt(s, IN(0.7), IN(5.05), IN(10.2), IN(1.0),
    [("Klaravex is an AI-first managed IT & security company for US small "
      "business — investor overview.", 17, False, MUTE)])
txt(s, IN(0.7), IN(6.55), IN(11.0), IN(0.4),
    [("CONFIDENTIAL · AUGUST 2026", 11, True, MUTE)])

# ---- 2 · hook ---------------------------------------------------------------
s = slide()
txt(s, IN(0.7), IN(1.5), IN(12.0), IN(2.6),
    [("89%", 150, True, INDIGO_L)])
txt(s, IN(0.75), IN(4.35), IN(11.4), IN(1.6),
    [("of IT support work no longer needs a human.", 30, True, WHITE),
     ("The industry still bills as if all of it does.", 30, False, MUTE)],
    space_before=8)
pagefoot(s, 2)

# ---- 3 · problem ------------------------------------------------------------
s = slide()
kicker(s, "The problem")
headline(s, "Small business IT is broken in three places")
probs = [
    ("Priced out", "Compliance-grade IT — HIPAA, SOC 2, ISO 27001 readiness — "
     "is enterprise-priced because it is delivered with enterprise headcount."),
    ("Misaligned", "Per-ticket billing rewards a full queue. Vendor commissions "
     "shape recommendations. The people grading the environment are the ones "
     "running it."),
    ("Low trust", "Consumers get truck-roll prices and urgency-driven scam "
     "“help.” Nobody in the market leads with honesty."),
]
cw = IN(3.95)
for i, (head, body) in enumerate(probs):
    x = IN(0.7) + (cw + IN(0.2)) * i
    card(s, x, IN(2.25), cw, IN(3.9))
    grad_bar(s, x + IN(0.3), IN(2.6), IN(0.45), IN(0.07))
    txt(s, x + IN(0.3), IN(2.8), cw - IN(0.6), IN(0.7),
        [(head, 21, True, WHITE)])
    txt(s, x + IN(0.3), IN(3.5), cw - IN(0.6), IN(2.4),
        [(body, 13.5, False, MUTE)])
pagefoot(s, 3)

# ---- 4 · why now ------------------------------------------------------------
s = slide()
kicker(s, "Why now")
headline(s, "The cost of resolution just collapsed")
rows = [
    ("AI can now close the loop", "Frontier models resolve real tickets end to "
     "end — triage, fix, documentation — not just deflect them to an FAQ."),
    ("Compliance pressure is rising", "Cyber-insurance and client audits now "
     "reach firms with 10 seats, not 1,000."),
    ("Incumbents cannot follow", "Labor-based MSPs would have to fire their own "
     "delivery model to match AI-first economics."),
]
for i, (head, body) in enumerate(rows):
    y = IN(2.3) + IN(1.45) * i
    card(s, IN(0.7), y, IN(11.9), IN(1.25))
    txt(s, IN(1.05), y + IN(0.18), IN(3.6), IN(0.95),
        [(head, 17, True, INDIGO_L)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, IN(4.9), y + IN(0.18), IN(7.4), IN(0.95),
        [(body, 14, False, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
pagefoot(s, 4)

# ---- 5 · product: two surfaces, one engine (diagram) ------------------------
s = slide()
kicker(s, "The product")
headline(s, "Two surfaces. One AI engine.")
# engine base
card(s, IN(2.4), IN(4.6), IN(8.5), IN(1.7), fill=INK2, line=INDIGO)
txt(s, IN(2.4), IN(4.78), IN(8.5), IN(0.5),
    [("KLARAVEX AI ENGINE", 15, True, INDIGO_L)], align=PP_ALIGN.CENTER)
txt(s, IN(2.7), IN(5.28), IN(7.9), IN(0.9),
    [("Voice · SMS · chat · email → one AI coordinator, one guardrail policy, "
      "one escalation path. Always labeled as AI.", 12.5, False, MUTE)],
    align=PP_ALIGN.CENTER)
# two surface cards
card(s, IN(1.15), IN(2.15), IN(5.2), IN(2.1), fill=INDIGO)
txt(s, IN(1.5), IN(2.35), IN(4.5), IN(0.5), [("klaravex.com", 19, True, WHITE)])
txt(s, IN(1.5), IN(2.9), IN(4.55), IN(1.25),
    [("B2B managed IT & security. Flat fee per user, tiers from $75 to $250. "
      "2-hour senior response, zero vendor commissions.", 12.5, False, WHITE)])
card(s, IN(6.95), IN(2.15), IN(5.2), IN(2.1), fill=VIOLET)
txt(s, IN(7.3), IN(2.35), IN(4.5), IN(0.5),
    [("personal.klaravex.com", 19, True, WHITE)])
txt(s, IN(7.3), IN(2.9), IN(4.55), IN(1.25),
    [("Consumer help by Klaravex AI. $29 sessions, plans from $29/month, "
      "free scam & hack recovery.", 12.5, False, WHITE)])
# connectors
for x in (IN(3.6), IN(9.4)):
    c = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, IN(4.25), IN(0.09), IN(0.36))
    c.fill.solid(); c.fill.fore_color.rgb = MUTE; c.line.fill.background()
    c.shadow.inherit = False
pagefoot(s, 5)

# ---- 6 · how it works (flow diagram) ---------------------------------------
s = slide()
kicker(s, "How it works")
headline(s, "Every issue takes the same honest path")
steps = [
    ("Any channel,\nany hour", "Phone, SMS, chat, email — 24/7, every time zone."),
    ("AI resolves\n89%", "Tier 1–2, monitoring, provisioning, onboarding, reporting — instantly."),
    ("Guardrails\non every turn", "Policy filters on input and output; escalation enforced at two boundaries."),
    ("Senior engineer\nfor the 11%", "Security, architecture, compliance — 2-hour commitment, service credit if missed."),
]
cw = IN(2.78)
for i, (head, body) in enumerate(steps):
    x = IN(0.7) + (cw + IN(0.32)) * i
    fill = INDIGO if i in (1,) else INK2
    card(s, x, IN(2.5), cw, IN(3.3), fill=fill)
    txt(s, x + IN(0.25), IN(2.75), cw - IN(0.5), IN(1.05),
        [(head, 16.5, True, WHITE)])
    txt(s, x + IN(0.25), IN(3.9), cw - IN(0.5), IN(1.7),
        [(body, 12, False, MUTE if fill == INK2 else WHITE)])
    if i < 3:
        arrow(s, x + cw + IN(0.02), IN(3.95), w=IN(0.29))
txt(s, IN(0.7), IN(6.1), IN(11.9), IN(0.5),
    [("Every action is logged; a human-review console owns approvals, billing "
      "and exceptions.", 12.5, False, MUTE)])
pagefoot(s, 6)

# ---- 7 · business model -----------------------------------------------------
s = slide()
kicker(s, "Business model")
headline(s, "Flat, recurring, conflict-free")
tiers = [
    ("Foundation", "from $75", "managed IT core", INK2),
    ("Assurance", "from $100", "+ security & identity", INK2),
    ("Directive", "from $150", "compliance readiness + MDR + vCISO — where "
     "enterprise deals start", INDIGO),
]
cw = IN(3.95)
for i, (name, price, desc, fill) in enumerate(tiers):
    x = IN(0.7) + (cw + IN(0.2)) * i
    card(s, x, IN(2.3), cw, IN(2.9), fill=fill)
    txt(s, x + IN(0.3), IN(2.55), cw - IN(0.6), IN(0.45),
        [(name.upper(), 13, True, INDIGO_L if fill == INK2 else WHITE)])
    txt(s, x + IN(0.3), IN(3.0), cw - IN(0.6), IN(0.9),
        [(price, 33, True, WHITE)])
    txt(s, x + IN(0.3), IN(3.95), cw - IN(0.6), IN(1.1),
        [(desc, 12.5, False, MUTE if fill == INK2 else WHITE)])
txt(s, IN(0.7), IN(5.5), IN(11.9), IN(0.9),
    [("Per user, per month. No per-ticket billing — nobody profits from a full "
      "queue. No vendor commissions — recommendations carry no hidden margin. "
      "Consumer line adds $29 sessions, $29/mo Solo, $39/mo Family.",
      13.5, False, MUTE)])
pagefoot(s, 7)

# ---- 8 · unit economics -----------------------------------------------------
s = slide()
kicker(s, "Unit economics")
headline(s, "AI resolution is gross margin")
stats = [("89%", "of volume at near-zero\nmarginal cost"),
         ("11%", "the only labor\nin the loop"),
         ("24/7", "coverage with no\nnight-shift payroll")]
cw = IN(3.95)
for i, (big, small) in enumerate(stats):
    x = IN(0.7) + (cw + IN(0.2)) * i
    card(s, x, IN(2.35), cw, IN(3.3))
    txt(s, x, IN(2.85), cw, IN(1.3), [(big, 58, True, INDIGO_L)],
        align=PP_ALIGN.CENTER)
    txt(s, x + IN(0.3), IN(4.35), cw - IN(0.6), IN(1.1),
        [(small, 14, False, MUTE)], align=PP_ALIGN.CENTER)
txt(s, IN(0.7), IN(6.05), IN(11.9), IN(0.6),
    [("Every point of AI resolution is margin a labor-based MSP has to buy "
      "with headcount.", 14, True, WHITE)])
pagefoot(s, 8)

# ---- 9 · moat ---------------------------------------------------------------
s = slide()
kicker(s, "The moat")
headline(s, "The company runs on its own product")
rows = [
    ("Autonomous engineering", "A closed reason–act–review–verify loop ships "
     "features against a versioned spec, with quality gates and audit trails."),
    ("Autonomous revenue", "Growth streams — leads, outreach, content — run on "
     "independent schedulers with human-gated publishing. Cadence survives "
     "failures elsewhere."),
    ("One-screen operations", "Delivery, pipeline, finances and the agent fleet "
     "on a single operator console with honest status. Today an internal edge; "
     "tomorrow a product."),
    ("Radical transparency", "AI always labeled. No fake reviews. Public agent "
     "experiments. Trust compounds in a low-trust market."),
]
for i, (head, body) in enumerate(rows):
    y = IN(2.25) + IN(1.08) * i
    grad_bar(s, IN(0.7), y + IN(0.12), IN(0.07), IN(0.75))
    txt(s, IN(1.0), y, IN(3.5), IN(0.95), [(head, 16, True, WHITE)],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, IN(4.7), y, IN(7.9), IN(0.95), [(body, 13, False, MUTE)],
        anchor=MSO_ANCHOR.MIDDLE)
pagefoot(s, 9)

# ---- 10 · go-to-market ------------------------------------------------------
s = slide()
kicker(s, "Go-to-market")
headline(s, "A warm wedge, then compounding funnels")
gtm = [
    ("1", "Founder-led wedge", "Years of hands-on engagements — foundations, "
     "legal aid, professional services. Free written assessments open doors."),
    ("2", "Compliance verticals", "Small law, accounting and medical practices "
     "— where the premium tier is bought, not sold."),
    ("3", "Autonomous funnel", "Growth streams generate leads and content on "
     "cadence; a public AI-agent scoreboard proves the governance."),
    ("4", "Consumer flywheel", "Free scam recovery earns word of mouth in the "
     "exact demographic that buys the Family plan."),
]
for i, (n, head, body) in enumerate(gtm):
    x = IN(0.7) + IN(6.1) * (i % 2)
    y = IN(2.3) + IN(1.95) * (i // 2)
    card(s, x, y, IN(5.85), IN(1.75))
    txt(s, x + IN(0.28), y + IN(0.22), IN(0.7), IN(1.0), [(n, 30, True, INDIGO_L)])
    txt(s, x + IN(1.05), y + IN(0.22), IN(4.5), IN(0.5), [(head, 15.5, True, WHITE)])
    txt(s, x + IN(1.05), y + IN(0.72), IN(4.55), IN(0.95), [(body, 11.5, False, MUTE)])
pagefoot(s, 10)

# ---- 11 · traction placeholder ---------------------------------------------
s = slide()
kicker(s, "Traction")
headline(s, "Numbers go here — and only real ones")
card(s, IN(0.7), IN(2.35), IN(11.9), IN(3.3), fill=INK2, line=VIOLET)
txt(s, IN(0.7), IN(3.0), IN(11.9), IN(0.8),
    [("[ MRR · pipeline · assessments delivered · consumer sessions ]",
      20, True, INDIGO_L)], align=PP_ALIGN.CENTER)
txt(s, IN(1.6), IN(3.9), IN(10.1), IN(1.4),
    [("Populate from live scorecards before any external send. This deck "
      "ships with placeholders on principle: the same no-fabrication policy "
      "the product is built on.", 14, False, MUTE)], align=PP_ALIGN.CENTER)
txt(s, IN(0.7), IN(5.95), IN(11.9), IN(0.8),
    [("Verifiable today: Secure Score 32% → 78% in under 60 days on a led "
      "engagement · full multi-site network replacement over one weekend.",
      13, False, WHITE)], align=PP_ALIGN.CENTER)
pagefoot(s, 11)

# ---- 12 · close / ask -------------------------------------------------------
s = slide()
grad_bar(s, 0, 0, W, IN(0.14))
s.shapes.add_picture(LOGO_LIGHT, IN(0.7), IN(1.0), height=IN(0.5))
txt(s, IN(0.7), IN(2.6), IN(11.9), IN(1.9),
    [("Honest, AI-first IT services —", 40, True, WHITE),
     ("built to scale without headcount.", 40, True, INDIGO_L)], space_before=2)
txt(s, IN(0.7), IN(5.0), IN(11.0), IN(1.2),
    [("Raising to scale the AI operations layer, cut growth streams to full "
      "autonomy, and productize the operator console.", 16, False, MUTE),
     ("hello@klaravex.com · (833) 990-2069 · klaravex.com", 16, True, WHITE)],
    space_before=14)
prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
