#!/usr/bin/env python3
"""Build the Klaravex pitch deck suite (v2, 2026-08-24).

Six decks per spec:
- 00-investor-overview: External, corporate voice, no vendor names
- 01-b2b-sales: B2B prospects (law/accounting/medical SMBs), CTA to klaravex.com
- 02-consumer: Personal/Klaravex partners/press, CTA to personal.klaravex.com
- 03-growth-os: Internal, mark INTERNAL
- 04-klaravex-os: Internal, distinguish KLARAVEX-OS from Founders OS
- 05-technical-architecture: Internal, ground in ARCHITECTURE.md

Design: indigo #4F46E5 / violet #7C3AED gradients, near-black #1C1C1A, warm neutrals #F5F3EE
Fonts: Syne (display) / Inter (body), fallback to Noto Sans
16:9 full-bleed, no white-with-a-stripe
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Brand colors
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
VIOLET = RGBColor(0x7C, 0x3A, 0xED)
DARK = RGBColor(0x1C, 0x1C, 0x1A)
WARM_NEUTRAL = RGBColor(0xF5, 0xF3, 0xEE)
DIM = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Slide dimensions (16:9)
SLIDE_WIDTH, SLIDE_HEIGHT = Inches(13.333), Inches(7.5)

OUTPUT_DIR = "/home/anthony/klaravex/decks-2026-08-24/v2"
BRAND_DIR = "/home/anthony/klaravex/brand/exports"


def new_deck():
    """Create a new presentation with correct dimensions."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_WIDTH, SLIDE_HEIGHT
    return prs


def blank_slide(prs):
    """Add a blank slide."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rectangle(slide, x, y, w, h, fill=None, line=None, radius=0):
    """Add a rectangle shape with optional fill, line, and rounding."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line
        shape.line.width = Pt(2)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, font="Inter"):
    """Add textbox with formatted runs (text, size, bold, color)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text_str, size, bold, color) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text_str
        f = r.font
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = color
        f.name = font
    return tb


def add_image(slide, x, y, w, h, img_path):
    """Add image from file."""
    return slide.shapes.add_picture(img_path, x, y, w, h)


def title_slide(prs, kicker, title, subtitle, is_internal=False, dark=True):
    """Create title slide with logo and brand colors."""
    slide = blank_slide(prs)
    
    # Background
    bg = add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, 
                       DARK if dark else WARM_NEUTRAL)
    
    # Brand accent strip
    add_rectangle(slide, 0, SLIDE_HEIGHT - Inches(0.18), SLIDE_WIDTH, Inches(0.18), INDIGO)
    
    # Logo (top right)
    logo_path = f"{BRAND_DIR}/klaravex-logo-light-2x.png" if dark else f"{BRAND_DIR}/klaravex-logo-dark-2x.png"
    add_image(slide, SLIDE_WIDTH - Inches(2.5), Inches(0.5), Inches(2), Inches(0.5), logo_path)
    
    # Text colors
    text_color = WHITE if dark else DARK
    subtitle_color = RGBColor(0xBB, 0xBB, 0xBB) if dark else DIM
    
    # Kicker
    add_text(slide, Inches(0.9), Inches(2.0), SLIDE_WIDTH - Inches(3), Inches(0.5),
             [(kicker.upper(), 14, True, INDIGO)])
    
    # Title
    add_text(slide, Inches(0.9), Inches(2.5), SLIDE_WIDTH - Inches(3), Inches(2.0),
             [(title, 46, True, text_color)], font="Syne")
    
    # Subtitle
    add_text(slide, Inches(0.9), Inches(4.4), SLIDE_WIDTH - Inches(3), Inches(1.6),
             [(subtitle, 18, False, subtitle_color)])
    
    # Internal watermark if needed
    if is_internal:
        add_text(slide, Inches(0.9), Inches(0.9), SLIDE_WIDTH - Inches(3), Inches(0.5),
                 [("INTERNAL", 12, True, INDIGO if dark else DIM)])
    
    return slide


def stat_slide(prs, title, stats, is_internal=False, footer=None):
    """Create a slide with big numeric stats in rounded cards."""
    slide = blank_slide(prs)
    
    # Background
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, WARM_NEUTRAL)
    
    # Brand strip
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, Inches(0.14), INDIGO)
    
    # Internal watermark
    if is_internal:
        add_text(slide, Inches(0.9), Inches(0.9), SLIDE_WIDTH - Inches(3), Inches(0.5),
                 [("INTERNAL", 12, True, DIM)])
    
    # Title
    add_text(slide, Inches(0.7), Inches(0.45), SLIDE_WIDTH - Inches(1.4), Inches(0.9),
             [(title, 30, True, DARK)], font="Syne")
    
    # Stats cards
    num_stats = len(stats)
    card_width = Inches(12.0 / num_stats)
    for i, (big, label) in enumerate(stats):
        x = Inches(0.7) + card_width * i
        card = add_rectangle(slide, x, Inches(2.0), card_width - Inches(0.3), Inches(3.4),
                           WARM_NEUTRAL)
        add_text(slide, x + Inches(0.2), Inches(2.6), card_width - Inches(0.7), Inches(1.2),
                 [(big, 44, True, INDIGO)], align=PP_ALIGN.CENTER, font="Syne")
        add_text(slide, x + Inches(0.2), Inches(3.9), card_width - Inches(0.7), Inches(1.3),
                 [(label, 14, False, DARK)], align=PP_ALIGN.CENTER)
    
    if footer:
        add_text(slide, Inches(0.75), Inches(6.5), SLIDE_WIDTH - Inches(1.5), Inches(0.7),
                 [(footer, 13, False, DIM)])
    
    return slide


def claim_slide(prs, claim, support, is_internal=False, footer=None):
    """Create a slide with one big claim and minimal supporting text."""
    slide = blank_slide(prs)
    
    # Background
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, WARM_NEUTRAL)
    
    # Brand strip
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, Inches(0.14), INDIGO)
    
    # Internal watermark
    if is_internal:
        add_text(slide, Inches(0.9), Inches(0.9), SLIDE_WIDTH - Inches(3), Inches(0.5),
                 [("INTERNAL", 12, True, DIM)])
    
    # Claim
    add_text(slide, Inches(0.7), Inches(1.5), SLIDE_WIDTH - Inches(1.4), Inches(1.5),
             [(claim, 36, True, DARK)], font="Syne")
    
    # Supporting text
    add_text(slide, Inches(0.7), Inches(3.2), SLIDE_WIDTH - Inches(1.4), Inches(2.5),
             [(support, 16, False, DIM)])
    
    if footer:
        add_text(slide, Inches(0.75), Inches(6.5), SLIDE_WIDTH - Inches(1.5), Inches(0.7),
                 [(footer, 13, False, DIM)])
    
    return slide


def diagram_slide(prs, title, diagram_type, is_internal=False):
    """Create a slide with a diagram (ticket flow, Growth OS, estate map)."""
    slide = blank_slide(prs)
    
    # Background
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, WARM_NEUTRAL)
    
    # Brand strip
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, Inches(0.14), INDIGO)
    
    # Internal watermark
    if is_internal:
        add_text(slide, Inches(0.9), Inches(0.9), SLIDE_WIDTH - Inches(3), Inches(0.5),
                 [("INTERNAL", 12, True, DIM)])
    
    # Title
    add_text(slide, Inches(0.7), Inches(0.45), SLIDE_WIDTH - Inches(1.4), Inches(0.9),
             [(title, 30, True, DARK)], font="Syne")
    
    # Diagram based on type
    diagram_area = add_rectangle(slide, Inches(0.7), Inches(1.5), SLIDE_WIDTH - Inches(1.4), Inches(5.0),
                               RGBColor(0xFF, 0xFF, 0xFF), line=INDIGO)
    
    if diagram_type == "Klaravex operational engine":
        # Simple 3-layer diagram
        layer1 = add_rectangle(slide, Inches(2.0), Inches(2.5), Inches(3.0), Inches(1.0), INDIGO)
        add_text(slide, Inches(2.2), Inches(2.7), Inches(2.6), Inches(0.6),
                 [("AI Engine", 16, True, WHITE)], align=PP_ALIGN.CENTER)
        
        layer2 = add_rectangle(slide, Inches(5.5), Inches(2.5), Inches(3.0), Inches(1.0), WARM_NEUTRAL, line=INDIGO)
        add_text(slide, Inches(5.7), Inches(2.7), Inches(2.6), Inches(0.6),
                 [("Operations", 16, True, INDIGO)], align=PP_ALIGN.CENTER)
        
        layer3 = add_rectangle(slide, Inches(9.0), Inches(2.5), Inches(3.0), Inches(1.0), WARM_NEUTRAL, line=INDIGO)
        add_text(slide, Inches(9.2), Inches(2.7), Inches(2.6), Inches(0.6),
                 [("Surfaces", 16, True, INDIGO)], align=PP_ALIGN.CENTER)
        
        # Arrows
        add_text(slide, Inches(4.9), Inches(2.8), Inches(0.5), Inches(0.4),
                 [("→", 20, True, DIM)], align=PP_ALIGN.CENTER)
        add_text(slide, Inches(8.4), Inches(2.8), Inches(0.5), Inches(0.4),
                 [("→", 20, True, DIM)], align=PP_ALIGN.CENTER)
    
    elif diagram_type == "Support ticket lifecycle":
        # Ticket flow diagram
        step1 = add_rectangle(slide, Inches(1.5), Inches(2.5), Inches(2.5), Inches(0.8), INDIGO)
        add_text(slide, Inches(1.7), Inches(2.6), Inches(2.1), Inches(0.6),
                 [("Ticket Created", 14, True, WHITE)], align=PP_ALIGN.CENTER)
        
        step2 = add_rectangle(slide, Inches(4.5), Inches(2.5), Inches(2.5), Inches(0.8), WARM_NEUTRAL, line=INDIGO)
        add_text(slide, Inches(4.7), Inches(2.6), Inches(2.1), Inches(0.6),
                 [("AI Triage", 14, True, INDIGO)], align=PP_ALIGN.CENTER)
        
        step3 = add_rectangle(slide, Inches(7.5), Inches(2.5), Inches(2.5), Inches(0.8), WARM_NEUTRAL, line=INDIGO)
        add_text(slide, Inches(7.7), Inches(2.6), Inches(2.1), Inches(0.6),
                 [("AI Resolution", 14, True, INDIGO)], align=PP_ALIGN.CENTER)
        
        step4 = add_rectangle(slide, Inches(10.5), Inches(2.5), Inches(2.5), Inches(0.8), WARM_NEUTRAL, line=INDIGO)
        add_text(slide, Inches(10.7), Inches(2.6), Inches(2.1), Inches(0.6),
                 [("Human Escalation", 14, True, INDIGO)], align=PP_ALIGN.CENTER)
        
        # Arrows
        add_text(slide, Inches(4.0), Inches(2.7), Inches(0.4), Inches(0.4),
                 [("→", 18, True, DIM)], align=PP_ALIGN.CENTER)
        add_text(slide, Inches(7.0), Inches(2.7), Inches(0.4), Inches(0.4),
                 [("→", 18, True, DIM)], align=PP_ALIGN.CENTER)
        add_text(slide, Inches(10.0), Inches(2.7), Inches(0.4), Inches(0.4),
                 [("→", 18, True, DIM)], align=PP_ALIGN.CENTER)
    
    elif diagram_type == "Consumer support session":
        # Simple session flow diagram
        step1 = add_rectangle(slide, Inches(3.0), Inches(2.5), Inches(2.5), Inches(0.8), INDIGO)
        add_text(slide, Inches(3.2), Inches(2.6), Inches(2.1), Inches(0.6),
                 [("Start Chat", 14, True, WHITE)], align=PP_ALIGN.CENTER)
        
        step2 = add_rectangle(slide, Inches(6.0), Inches(2.5), Inches(2.5), Inches(0.8), WARM_NEUTRAL, line=INDIGO)
        add_text(slide, Inches(6.2), Inches(2.6), Inches(2.1), Inches(0.6),
                 [("Describe Problem", 14, True, INDIGO)], align=PP_ALIGN.CENTER)
        
        step3 = add_rectangle(slide, Inches(9.0), Inches(2.5), Inches(2.5), Inches(0.8), WARM_NEUTRAL, line=INDIGO)
        add_text(slide, Inches(9.2), Inches(2.6), Inches(2.1), Inches(0.6),
                 [("Guided Fix", 14, True, INDIGO)], align=PP_ALIGN.CENTER)
        
        # Arrows
        add_text(slide, Inches(5.4), Inches(2.7), Inches(0.4), Inches(0.4),
                 [("→", 18, True, DIM)], align=PP_ALIGN.CENTER)
        add_text(slide, Inches(8.4), Inches(2.7), Inches(0.4), Inches(0.4),
                 [("→", 18, True, DIM)], align=PP_ALIGN.CENTER)
    
    elif diagram_type == "Four-layer architecture":
        # Growth OS layers diagram
        layer1 = add_rectangle(slide, Inches(5.0), Inches(2.5), Inches(3.5), Inches(0.8), INDIGO)
        add_text(slide, Inches(5.2), Inches(2.6), Inches(3.1), Inches(0.6),
                 [("Revenue Agents", 14, True, WHITE)], align=PP_ALIGN.CENTER)
        
        layer2 = add_rectangle(slide, Inches(5.0), Inches(3.5), Inches(3.5), Inches(0.8), WARM_NEUTRAL, line=INDIGO)
        add_text(slide, Inches(5.2), Inches(3.6), Inches(3.1), Inches(0.6),
                 [("Ops Glue (n8n)", 14, True, INDIGO)], align=PP_ALIGN.CENTER)
        
        layer3 = add_rectangle(slide, Inches(5.0), Inches(4.5), Inches(3.5), Inches(0.8), WARM_NEUTRAL, line=INDIGO)
        add_text(slide, Inches(5.2), Inches(4.6), Inches(3.1), Inches(0.6),
                 [("Growth API", 14, True, INDIGO)], align=PP_ALIGN.CENTER)
        
        layer4 = add_rectangle(slide, Inches(5.0), Inches(5.5), Inches(3.5), Inches(0.8), WARM_NEUTRAL, line=INDIGO)
        add_text(slide, Inches(5.2), Inches(5.6), Inches(3.1), Inches(0.6),
                 [("KLARAVEX-OS", 14, True, INDIGO)], align=PP_ALIGN.CENTER)
        
        # Arrows
        add_text(slide, Inches(6.5), Inches(3.3), Inches(0.4), Inches(0.4),
                 [("→", 18, True, DIM)], align=PP_ALIGN.CENTER)
        add_text(slide, Inches(6.5), Inches(4.3), Inches(0.4), Inches(0.4),
                 [("→", 18, True, DIM)], align=PP_ALIGN.CENTER)
        add_text(slide, Inches(6.5), Inches(5.3), Inches(0.4), Inches(0.4),
                 [("→", 18, True, DIM)], align=PP_ALIGN.CENTER)
    
    elif diagram_type == "Klaravex platform estate":
        # Estate map diagram
        legacy = add_rectangle(slide, Inches(2.0), Inches(3.5), Inches(2.5), Inches(0.8), RGBColor(0x99, 0x99, 0x99))
        add_text(slide, Inches(2.2), Inches(3.6), Inches(2.1), Inches(0.6),
                 [("Klaravex (Legacy)", 12, True, WHITE)], align=PP_ALIGN.CENTER)
        
        transition = add_rectangle(slide, Inches(5.0), Inches(3.5), Inches(2.5), Inches(0.8), INDIGO)
        add_text(slide, Inches(5.2), Inches(3.6), Inches(2.1), Inches(0.6),
                 [("Klaravex2.0", 12, True, WHITE)], align=PP_ALIGN.CENTER)
        
        modern = add_rectangle(slide, Inches(8.0), Inches(3.5), Inches(2.5), Inches(0.8), VIOLET)
        add_text(slide, Inches(8.2), Inches(3.6), Inches(2.1), Inches(0.6),
                 [("KLARAVEX-OS", 12, True, WHITE)], align=PP_ALIGN.CENTER)
        
        # Arrows
        add_text(slide, Inches(4.4), Inches(3.7), Inches(0.5), Inches(0.4),
                 [("→", 18, True, DIM)], align=PP_ALIGN.CENTER)
        add_text(slide, Inches(7.4), Inches(3.7), Inches(0.5), Inches(0.4),
                 [("→", 18, True, DIM)], align=PP_ALIGN.CENTER)
    
    # Diagram label
    add_text(slide, Inches(0.7), Inches(6.6), SLIDE_WIDTH - Inches(1.4), Inches(0.7),
             [(f"{diagram_type} diagram", 13, False, DIM)], align=PP_ALIGN.CENTER)
    
    return slide


def tiers_slide(prs, title, tiers, is_internal=False):
    """Create a slide with service tiers in colored cards."""
    slide = blank_slide(prs)
    
    # Background
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, WARM_NEUTRAL)
    
    # Brand strip
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, Inches(0.14), INDIGO)
    
    # Internal watermark
    if is_internal:
        add_text(slide, Inches(0.9), Inches(0.9), SLIDE_WIDTH - Inches(3), Inches(0.5),
                 [("INTERNAL", 12, True, DIM)])
    
    # Title
    add_text(slide, Inches(0.7), Inches(0.45), SLIDE_WIDTH - Inches(1.4), Inches(0.9),
             [(title, 30, True, DARK)], font="Syne")
    
    # Tier cards
    num_tiers = len(tiers)
    card_width = Inches(12.0 / num_tiers)
    for i, (name, price, desc) in enumerate(tiers):
        x = Inches(0.7) + card_width * i
        card = add_rectangle(slide, x, Inches(1.5), card_width - Inches(0.3), Inches(4.5),
                           INDIGO if i == num_tiers - 1 else WARM_NEUTRAL)
        text_color = WHITE if i == num_tiers - 1 else DARK
        price_color = INDIGO if i != num_tiers - 1 else WHITE
        
        add_text(slide, x + Inches(0.15), Inches(1.8), card_width - Inches(0.6), Inches(0.8),
                 [(name, 18, True, text_color)], font="Syne")
        add_text(slide, x + Inches(0.15), Inches(2.6), card_width - Inches(0.6), Inches(0.6),
                 [(price, 24, True, price_color)], font="Syne")
        add_text(slide, x + Inches(0.15), Inches(3.3), card_width - Inches(0.6), Inches(2.0),
                 [(desc, 13, False, text_color)])
    
    return slide


def cta_slide(prs, title, lines, url, is_internal=False, dark=True):
    """Create a call-to-action slide with strong color contrast."""
    slide = blank_slide(prs)
    
    # Background
    bg = add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT,
                       DARK if dark else WARM_NEUTRAL)
    
    # Brand accent strip
    add_rectangle(slide, 0, SLIDE_HEIGHT - Inches(0.18), SLIDE_WIDTH, Inches(0.18), INDIGO)
    
    # Logo (top right)
    logo_path = f"{BRAND_DIR}/klaravex-logo-light-2x.png" if dark else f"{BRAND_DIR}/klaravex-logo-dark-2x.png"
    add_image(slide, SLIDE_WIDTH - Inches(2.5), Inches(0.5), Inches(2), Inches(0.5), logo_path)
    
    # Text colors
    text_color = WHITE if dark else DARK
    desc_color = RGBColor(0xCC, 0xCC, 0xCC) if dark else DIM
    url_color = RGBColor(0x81, 0x8C, 0xF8) if dark else INDIGO
    
    # Title
    add_text(slide, Inches(0.9), Inches(2.3), SLIDE_WIDTH - Inches(3), Inches(1.2),
             [(title, 40, True, text_color)], font="Syne")
    
    # Description lines
    y = 3.7
    for line in lines:
        add_text(slide, Inches(0.9), Inches(y), SLIDE_WIDTH - Inches(3), Inches(0.5),
                 [(line, 17, False, desc_color)])
        y += 0.55
    
    # URL
    add_text(slide, Inches(0.9), Inches(y + 0.3), SLIDE_WIDTH - Inches(3), Inches(0.7),
             [(url, 24, True, url_color)])
    
    # Internal watermark if needed
    if is_internal:
        add_text(slide, Inches(0.9), Inches(0.9), SLIDE_WIDTH - Inches(3), Inches(0.5),
                 [("INTERNAL", 12, True, INDIGO if dark else DIM)])
    
    return slide


# ---------------------------------------------------------------------------------
# Deck 00: Investor Overview
# ---------------------------------------------------------------------------------
def build_investor_deck():
    prs = new_deck()
    
    title_slide(prs, "Klaravex LLC — Investor Deck",
                "AI-first managed IT & security for US businesses",
                "89% of issues resolved by AI instantly, 24/7. The 11% that need judgment get a senior engineer in under two hours.",
                is_internal=False, dark=True)
    
    claim_slide(prs, "The problem with IT support",
                "Small US businesses pay $150–300/hr for reactive IT that overpromises and underdelivers. Security gaps cost them millions in breaches.",
                is_internal=False)
    
    claim_slide(prs, "Why now",
                "AI can automate 89% of routine IT work instantly. Senior engineers focus on the 11% that matters — security, architecture, compliance.",
                is_internal=False)
    
    claim_slide(prs, "One engine, two surfaces",
                "Klaravex runs on a single AI-powered operations platform serving both B2B businesses and B2C consumers.",
                is_internal=False)
    
    diagram_slide(prs, "How it works", "Klaravex operational engine", is_internal=False)
    
    stat_slide(prs, "Unit economics", [
        ("89%", "of issues resolved by AI instantly"),
        ("11%", "handled by senior engineers"),
        ("<2hrs", "engineer response time"),
        ("$75–250", "per user per month (B2B)"),
    ], is_internal=False)
    
    claim_slide(prs, "Moat: The company runs on its own product",
                "Klaravex's entire operations, sales, and support run on the same platform it sells — ensuring constant product-market fit.",
                is_internal=False)
    
    claim_slide(prs, "Go-to-market",
                "Targeting law, accounting, and medical SMBs with free security assessments and transparent pricing.",
                is_internal=False)
    
    claim_slide(prs, "Traction",
                "[PLACEHOLDER — Pull real numbers from scorecards. No fabrication.]",
                is_internal=False)
    
    claim_slide(prs, "The team",
                "Founder with 15+ years in enterprise IT and security. Senior engineering team with cloud and AI expertise.",
                is_internal=False)
    
    cta_slide(prs, "Invest in the future of IT support",
              ["AI-first, transparent, and accountable. Built for US businesses."],
              "klaravex.com", is_internal=False, dark=True)
    
    prs.save(f"{OUTPUT_DIR}/00-investor-overview.pptx")
    print("Generated 00-investor-overview.pptx")


# ---------------------------------------------------------------------------------
# Deck 01: B2B Sales
# ---------------------------------------------------------------------------------
def build_b2b_sales_deck():
    prs = new_deck()
    
    title_slide(prs, "Klaravex — B2B Sales Deck",
                "AI-first managed IT & security for US businesses",
                "89% of issues resolved instantly, 24/7. Senior engineers in under 2 hours. Transparent pricing.",
                is_internal=False, dark=True)
    
    claim_slide(prs, "Your IT is a liability",
                "Small law, accounting, and medical firms face constant security risks and expensive reactive support.",
                is_internal=False)
    
    claim_slide(prs, "What you get",
                "AI-powered support, security hardening, compliance readiness, and a senior engineer on call.",
                is_internal=False)
    
    diagram_slide(prs, "Ticket flow", "Support ticket lifecycle", is_internal=False)
    
    tiers_slide(prs, "Service tiers", [
        ("Foundation", "$75–100 / user / mo", "AI-first support, patching, monitoring, backup, network management"),
        ("Assurance", "$100–150 / user / mo", "Everything in Foundation plus security hardening and identity management"),
        ("Directive", "$150–250 / user / mo", "Compliance readiness (HIPAA/SOC 2/ISO 27001) and vCISO support"),
    ], is_internal=False)
    
    stat_slide(prs, "Proof", [
        ("32→78", "Secure Score improvement in <60 days"),
        ("1 weekend", "Full network replacement"),
        ("2 hrs", "Engineer response commitment"),
        ("$0", "Vendor commissions"),
    ], is_internal=False)
    
    cta_slide(prs, "Ready for better IT?",
              ["Free assessment: 45 minutes, written report, no strings attached."],
              "klaravex.com", is_internal=False, dark=True)
    
    prs.save(f"{OUTPUT_DIR}/01-b2b-sales.pptx")
    print("Generated 01-b2b-sales.pptx")


# ---------------------------------------------------------------------------------
# Deck 02: Consumer
# ---------------------------------------------------------------------------------
def build_consumer_deck():
    prs = new_deck()
    
    title_slide(prs, "Klaravex Personal — Consumer Deck",
                "Plain-English tech help, any hour",
                "AI-powered support labeled as AI. $29 sessions, monthly plans, free scam recovery.",
                is_internal=False, dark=True)
    
    stat_slide(prs, "Simple pricing", [
        ("$29", "One-off session"),
        ("$29/mo", "Solo plan"),
        ("$39/mo", "Family plan"),
        ("Free", "Scam & hack recovery"),
    ], is_internal=False)
    
    claim_slide(prs, "Who it serves",
                "Parents, grandparents, scam victims, job seekers — anyone who hates tech jargon.",
                is_internal=False)
    
    diagram_slide(prs, "Session flow", "Consumer support session", is_internal=False)
    
    claim_slide(prs, "Trust by design",
                "Always labeled as AI. No fake humans, no upsell scripts, transparent pricing.",
                is_internal=False)
    
    cta_slide(prs, "Tech that works for you",
              ["A session is $29. Scam or hack? Help is free. Chat any hour."],
              "personal.klaravex.com", is_internal=False, dark=True)
    
    prs.save(f"{OUTPUT_DIR}/02-consumer.pptx")
    print("Generated 02-consumer.pptx")


# ---------------------------------------------------------------------------------
# Deck 03: Growth OS (Internal)
# ---------------------------------------------------------------------------------
def build_growth_os_deck():
    prs = new_deck()
    
    title_slide(prs, "INTERNAL — Klaravex Growth OS",
                "The strangler-fig rebuild",
                "A parallel product tree that takes over revenue generation stream by stream while the monolith serves.",
                is_internal=True, dark=True)
    
    claim_slide(prs, "Why a second tree",
                "The legacy monolith works but limits growth cadence. Growth OS runs independently on systemd timers.",
                is_internal=True)
    
    diagram_slide(prs, "Growth OS layers", "Four-layer architecture", is_internal=True)
    
    claim_slide(prs, "Migration strategy",
                "Shadow → cutover per stream. Live monolith remains production until each stream is verified.",
                is_internal=True)
    
    cta_slide(prs, "Growth OS",
              ["Operator access: KLARAVEX-OS cockpit on :4100"],
              "internal.klaravex.com", is_internal=True, dark=True)
    
    prs.save(f"{OUTPUT_DIR}/03-growth-os.pptx")
    print("Generated 03-growth-os.pptx")


# ---------------------------------------------------------------------------------
# Deck 04: Klaravex OS (Internal)
# ---------------------------------------------------------------------------------
def build_klaravex_os_deck():
    prs = new_deck()
    
    title_slide(prs, "INTERNAL — KLARAVEX-OS",
                "The operator cockpit",
                "Distinct from Founders OS (client portal). Runs on :4100 for internal operations.",
                is_internal=True, dark=True)
    
    claim_slide(prs, "KLARAVEX-OS vs Founders OS",
                "KLARAVEX-OS (:4100) = internal operations. Founders OS = client portal for end-users.",
                is_internal=True)
    
    claim_slide(prs, "Core features",
                "Real-time monitoring, agent orchestration, billing management, and compliance dashboards.",
                is_internal=True)
    
    claim_slide(prs, "Architecture",
                "Built on FastAPI, PostgreSQL, and systemd timers. Calls Growth API exclusively.",
                is_internal=True)
    
    cta_slide(prs, "Access KLARAVEX-OS",
              ["Local: http://localhost:4100 | Production: [internal URL]"],
              "localhost:4100", is_internal=True, dark=True)
    
    prs.save(f"{OUTPUT_DIR}/04-klaravex-os.pptx")
    print("Generated 04-klaravex-os.pptx")


# ---------------------------------------------------------------------------------
# Deck 05: Technical Architecture (Internal)
# ---------------------------------------------------------------------------------
def build_technical_architecture_deck():
    prs = new_deck()
    
    title_slide(prs, "INTERNAL — Technical Architecture",
                "Klaravex platform architecture",
                "Ground truth in ARCHITECTURE.md, Klaravex2.0/MIGRATION.md, and klaravex-os README.",
                is_internal=True, dark=True)
    
    claim_slide(prs, "Estate map",
                "Klaravex → Klaravex2.0 → KLARAVEX-OS: incremental migration from legacy to modern stack.",
                is_internal=True)
    
    diagram_slide(prs, "Estate map", "Klaravex platform estate", is_internal=True)
    
    claim_slide(prs, "Growth OS architecture",
                "Four layers: Revenue agents → Ops glue (n8n) → Growth API → KLARAVEX-OS.",
                is_internal=True)
    
    claim_slide(prs, "Core technologies",
                "Python, FastAPI, PostgreSQL, systemd timers, Claude AI agents, UniFi network management.",
                is_internal=True)
    
    claim_slide(prs, "Security boundaries",
                "Zero-trust networking, conditional access, encrypted messaging, audit trails per compliance requirements.",
                is_internal=True)
    
    cta_slide(prs, "Documentation",
              ["ARCHITECTURE.md · Klaravex2.0/MIGRATION.md · klaravex-os README"],
              "internal-docs.klaravex.com", is_internal=True, dark=True)
    
    prs.save(f"{OUTPUT_DIR}/05-technical-architecture.pptx")
    print("Generated 05-technical-architecture.pptx")


if __name__ == "__main__":
    print("Building Klaravex pitch deck suite v2...")
    
    build_investor_deck()
    build_b2b_sales_deck()
    build_consumer_deck()
    build_growth_os_deck()
    build_klaravex_os_deck()
    build_technical_architecture_deck()
    
    print(f"All decks generated in {OUTPUT_DIR}")
    
    # Verify generation and render to PDF for checking
    import os
    import subprocess
    
    print("\nRendering PDFs for verification (requires LibreOffice)...")
    success = True
    
    for filename in os.listdir(OUTPUT_DIR):
        if filename.endswith(".pptx"):
            pptx_path = os.path.join(OUTPUT_DIR, filename)
            pdf_path = os.path.join(OUTPUT_DIR, filename.replace(".pptx", ".pdf"))
            
            try:
                subprocess.run([
                    "soffice", "--headless", "--convert-to", "pdf",
                    "--outdir", OUTPUT_DIR, pptx_path
                ], check=True, capture_output=True, text=True)
                print(f"  ✓ {filename} → {os.path.basename(pdf_path)}")
            except subprocess.CalledProcessError as e:
                print(f"  ✗ Failed to convert {filename}: {e.stderr}")
                success = False
    
    if success:
        print("\n✅ All decks rendered successfully. Check PDFs in v2/ directory.")
    else:
        print("\n❌ Some decks failed to render. Check errors above.")
