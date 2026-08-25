#!/usr/bin/env python3
"""Build the four Klaravex pitch decks (2026-08-24).

Decks 1-2 are external-safe (corporate voice: no internal codenames, no
infrastructure vendor names). Decks 3-4 are marked INTERNAL.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

INDIGO = RGBColor(0x4F, 0x46, 0xE5)
VIOLET = RGBColor(0x7C, 0x3A, 0xED)
DARK = RGBColor(0x1C, 0x1C, 0x1A)
DIM = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x10, 0xB9, 0x81)

W, H = Inches(13.333), Inches(7.5)


def new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def box(slide, x, y, w, h, fill=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.line.fill.background()
    if fill is not None:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    return sh


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT):
    """runs: list of (text, size, bold, color) tuples -> one paragraph each."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (t, size, bold, color) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = t
        f = r.font
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = color
        f.name = "Calibri"
    return tb


def title_slide(prs, kicker, title, sub, dark=False):
    s = blank(prs)
    bg = box(s, 0, 0, W, H, DARK if dark else WHITE)
    box(s, 0, H - Inches(0.18), W, Inches(0.18), INDIGO)
    fg = WHITE if dark else DARK
    text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.5),
         [(kicker.upper(), 14, True, INDIGO)])
    text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.8),
         [(title, 46, True, fg)])
    text(s, Inches(0.9), Inches(4.4), Inches(11.0), Inches(1.6),
         [(sub, 18, False, DIM if not dark else RGBColor(0xBB, 0xBB, 0xBB))])
    return s


def bullets_slide(prs, title, items, footer=None):
    """items: list of (head, body) or plain strings."""
    s = blank(prs)
    box(s, 0, 0, W, Inches(0.14), INDIGO)
    text(s, Inches(0.7), Inches(0.45), Inches(12.0), Inches(0.9),
         [(title, 30, True, DARK)])
    tb = s.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.9), Inches(5.4))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        head, body = it if isinstance(it, tuple) else (None, it)
        if head:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            r = p.add_run(); r.text = "  " + head
            r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = INDIGO
            p.space_before = Pt(10)
            if body:
                p2 = tf.add_paragraph()
                r2 = p2.add_run(); r2.text = "      " + body
                r2.font.size = Pt(14); r2.font.color.rgb = DARK
        else:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            r = p.add_run(); r.text = "  •  " + body
            r.font.size = Pt(16); r.font.color.rgb = DARK
            p.space_before = Pt(8)
    if footer:
        text(s, Inches(0.75), Inches(6.9), Inches(12.0), Inches(0.45),
             [(footer, 12, False, DIM)])
    return s


def stat_slide(prs, title, stats, footer=None):
    """stats: list of (big, label)."""
    s = blank(prs)
    box(s, 0, 0, W, Inches(0.14), INDIGO)
    text(s, Inches(0.7), Inches(0.45), Inches(12.0), Inches(0.9),
         [(title, 30, True, DARK)])
    n = len(stats)
    cw = Inches(12.0 / n)
    for i, (big, label) in enumerate(stats):
        x = Inches(0.7) + cw * i
        card = box(s, x, Inches(2.0), cw - Inches(0.3), Inches(3.4),
                   RGBColor(0xF5, 0xF3, 0xEE))
        text(s, x + Inches(0.2), Inches(2.6), cw - Inches(0.7), Inches(1.2),
             [(big, 44, True, INDIGO)], align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.2), Inches(3.9), cw - Inches(0.7), Inches(1.3),
             [(label, 14, False, DARK)], align=PP_ALIGN.CENTER)
    if footer:
        text(s, Inches(0.75), Inches(6.5), Inches(12.0), Inches(0.7),
             [(footer, 13, False, DIM)])
    return s


def flow_slide(prs, title, steps, footer=None):
    """steps: list of (head, body) rendered as a left-to-right chain."""
    s = blank(prs)
    box(s, 0, 0, W, Inches(0.14), INDIGO)
    text(s, Inches(0.7), Inches(0.45), Inches(12.0), Inches(0.9),
         [(title, 30, True, DARK)])
    n = len(steps)
    cw = Inches(12.2 / n)
    for i, (head, body) in enumerate(steps):
        x = Inches(0.6) + cw * i
        box(s, x, Inches(2.1), cw - Inches(0.35), Inches(3.2),
            INDIGO if i == 0 else RGBColor(0xF5, 0xF3, 0xEE))
        fg = WHITE if i == 0 else DARK
        text(s, x + Inches(0.15), Inches(2.3), cw - Inches(0.65), Inches(0.9),
             [(head, 16, True, fg if i == 0 else INDIGO)])
        text(s, x + Inches(0.15), Inches(3.1), cw - Inches(0.65), Inches(2.0),
             [(body, 12, False, fg)])
        if i < n - 1:
            text(s, x + cw - Inches(0.38), Inches(3.3), Inches(0.4), Inches(0.5),
                 [("→", 22, True, DIM)])
    if footer:
        text(s, Inches(0.75), Inches(6.6), Inches(12.0), Inches(0.6),
             [(footer, 13, False, DIM)])
    return s


def cta_slide(prs, title, lines, url):
    s = blank(prs)
    box(s, 0, 0, W, H, DARK)
    text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.2),
         [(title, 40, True, WHITE)])
    y = 3.7
    for ln in lines:
        text(s, Inches(0.9), Inches(y), Inches(11.5), Inches(0.5),
             [(ln, 17, False, RGBColor(0xCC, 0xCC, 0xCC))])
        y += 0.55
    text(s, Inches(0.9), Inches(y + 0.3), Inches(11.5), Inches(0.7),
         [(url, 24, True, RGBColor(0x81, 0x8C, 0xF8))])
    return s


OUT = "/home/anthony/klaravex/decks-2026-08-24/"

# ---------------------------------------------------------------- Deck 1: B2B
prs = new_deck()
title_slide(prs, "Klaravex LLC — pitch deck",
            "AI-first managed IT & security for US businesses",
            "89% of issues resolved by AI instantly, 24/7. The 11% that need judgment "
            "get a senior engineer in under two hours. Flat fee per user. "
            "No vendor commissions — ever.", dark=True)
stat_slide(prs, "The model in four numbers", [
    ("89%", "of issues resolved by AI instantly — Tier 1 & 2, monitoring, provisioning, onboarding, reporting"),
    ("11%", "handled personally by a senior engineer — security incidents, architecture, compliance, vendor decisions"),
    ("2 hrs", "senior engineer response commitment, backed by a service credit"),
    ("$0", "vendor commissions. Recommendations carry no hidden margin"),
])
bullets_slide(prs, "Who it is for", [
    ("Small law firms", "Client-file confidentiality, litigation-hold readiness, email security that holds up in discovery."),
    ("Accounting practices", "SOC 2 readiness, financial-data controls, tax-season uptime."),
    ("Medical offices", "Strict HIPAA readiness — access controls, audit trails, encrypted messaging, BAA discipline."),
    ("US SMBs on Microsoft 365 / Google Workspace / AWS", "Hardening, identity, Conditional Access, backup coverage, network infrastructure managed end to end."),
], footer="UniFi network infrastructure management is bundled into every tier.")
flow_slide(prs, "How a ticket flows", [
    ("Any channel, any hour", "Phone, SMS, web chat, and email all land in one AI support coordinator — Klaravex AI, always labeled as AI."),
    ("AI triage & resolution", "The coordinator scores the issue against six specialist engineering pillars and resolves routine work instantly."),
    ("Guardrails", "Every AI input and output passes policy filters. Escalation paths are enforced at two independent boundaries."),
    ("Human judgment", "Security incidents, architecture, and compliance route to a senior engineer with a 2-hour commitment."),
    ("Accountable record", "Every action logged; a human-review console covers approvals, billing, and exceptions."),
])
bullets_slide(prs, "Service tiers", [
    ("Foundation — ~$75–100 / user / month", "Managed IT core: AI-first support, patching, monitoring, backup coverage, network management."),
    ("Assurance — ~$100–150 / user / month", "Everything in Foundation plus security hardening, identity management, and compliance groundwork."),
    ("Directive — ~$150–250 / user / month", "Compliance readiness (HIPAA / SOC 2 / ISO 27001), managed detection & response, and vCISO support. Where enterprise conversations start."),
], footer="Flat fee per user. No per-ticket billing — nobody has an incentive to keep the queue full.")
bullets_slide(prs, "Proof, not promises", [
    ("Secure Score 32% → 78% in under 60 days", "An engagement our founder led: M365 hardening, MFA enforcement, Conditional Access, email security — left cyber-insurance ready."),
    ("Full network replacement over one weekend", "Firewalls and switching across a multi-site SMB. Staff noticed faster Wi-Fi on Monday and nothing else."),
    ("Two-hour response, in writing", "Missed commitments earn service credits automatically."),
])
cta_slide(prs, "Ready for IT that reports honestly?",
          ["Free assessment: 45 minutes, a written report, nothing attached.",
           "hello@klaravex.com · (833) 990-2069"],
          "klaravex.com")
prs.save(OUT + "01-klaravex-b2b.pptx")

# ------------------------------------------------------- Deck 2: Personal
prs = new_deck()
title_slide(prs, "Klaravex Personal — pitch deck",
            "Plain-English tech help, any hour",
            "Remote-only consumer tech support delivered by Klaravex AI — always labeled "
            "as AI, never pretending to be human. AI does the heavy lifting so a session "
            "stays $29 instead of a house call.", dark=True)
stat_slide(prs, "Pricing a person can say yes to", [
    ("$29", "one-off session — pay only when you need help"),
    ("$29/mo", "Solo plan — ongoing help for one person"),
    ("$39/mo", "Family plan — parents and grandparents covered"),
    ("Free", "scam & hack recovery — for anyone, no purchase required"),
], footer="Remote only — no house calls. That is how a session stays $29 instead of a truck roll.")
bullets_slide(prs, "Who it serves", [
    ("Parents and grandparents", "Patient, jargon-free help with the printer, the Wi-Fi, the phone, the scam text."),
    ("Scam and hack victims", "Free recovery help — accounts, passwords, damage control — because that should never be paywalled."),
    ("Job seekers and solo founders", "Resume & job-hunt tech kit, solo-business launch kit, AI skills coaching."),
    ("Anyone allergic to jargon", "Plain English, no judgment, any hour."),
])
flow_slide(prs, "How a session works", [
    ("Open the chat", "Klara — the Klaravex AI coordinator — greets you on personal.klaravex.com. Clearly labeled AI, from the first word."),
    ("Describe the problem", "In plain English. Klara remembers the conversation — no repeating yourself."),
    ("Guided fix", "Step-by-step, at your pace, with a screen-share when it helps."),
    ("Honest next steps", "If it needs more than a session, Klara says so in plain English — you decide."),
])
bullets_slide(prs, "Trust is the product", [
    ("Always labeled as AI", "Every chat is named Klaravex AI. No fake humans, no fake reviews, no invented ratings."),
    ("Privacy by design", "Consent asked up front; data handled per a plain-language privacy policy."),
    ("Transparent pricing", "The price on the page is the price. No upsell scripts."),
    ("Free where it matters most", "Scam and hack recovery is free because urgency should not be leverage."),
])
cta_slide(prs, "Tech that actually works for you",
          ["A session is $29 — monthly plans from $29.",
           "Scam or hack? Help is free. Start a chat any hour."],
          "personal.klaravex.com")
prs.save(OUT + "02-klaravex-personal.pptx")

# ---------------------------------------------------- Deck 3: Growth OS (2.0)
prs = new_deck()
title_slide(prs, "INTERNAL — Klaravex 2.0",
            "Growth OS: the strangler-fig rebuild",
            "A parallel product tree that takes over revenue generation stream by stream "
            "— while the production monolith keeps serving. Growth must keep running even "
            "if the legacy schedulers die.", dark=True)
bullets_slide(prs, "Why a second tree", [
    ("The monolith works — and that is the problem", "Growth cadence, agent charters, and product runtime share one crash domain and one scheduler (Celery beat)."),
    ("Strangler-fig, not rewrite", "Klaravex2.0 owns the new implementation home; live klaravex stays production until each stream cuts over."),
    ("Independence as a requirement", "Growth streams run on systemd timers — if Celery beat or the ops layer dies, revenue generation keeps its cadence."),
])
flow_slide(prs, "Four layers", [
    ("A — Revenue agents", "Charters + local outbox in Klaravex2.0/revenue-agents. Source of truth for agent behavior."),
    ("B — Ops glue (n8n)", "Optional callers of the Growth API. Never owns rubrics."),
    ("C — Growth API", "FastAPI control plane on :4210 + systemd timers + charter executor running Claude in background threads."),
    ("D — KLARAVEX-OS", "Operator cockpit on :4100. Calls the Growth API exclusively for growth — see deck 4."),
])
bullets_slide(prs, "Migration: shadow → cutover, per stream", [
    ("Phase 0–1", "Inventory legacy beat schedules; stand up the Growth API stub with auth and health checks."),
    ("Phase 2 — where we are", "Shadow: timers live and firing into the Growth API; legacy beat still owns truth; outboxes compared."),
    ("Phase 3", "Per-stream cutover: leads first (low blast radius), then freelance/ads, then the gated publish path (socials → seo-blog → kb → backlinks), gatekeeper last."),
    ("Phase 4", "Beat-kill test: stop the legacy scheduler in a controlled window and prove Growth cadence survives."),
    ("Phase 5", "Adapters (Clay / Taplio / outreach / WordPress) and full Layer-D wiring."),
], footer="Rollback per stream: disable the timer, re-enable the legacy schedule, keep the new outbox as forensic replay.")
bullets_slide(prs, "Guardrails carried over", [
    ("Human gate on publish", "Gated streams draft to the outbox; nothing publishes without the gatekeeper."),
    ("Scorecards", "Every stream reports cadence and output to a scorecard — silent failure is treated as failure."),
    ("Non-goals stay non-goals", "No replacement of the support runtime or deploys; no shared crash domains; no defense/CMMC surfaces."),
])
cta_slide(prs, "Status: Phase 2 shadow, timers live",
          ["Next milestone: first stream cutover (leads).",
           "Runbooks: MIGRATION.md · docs/cutover-checklist.md"],
          "Klaravex2.0 — internal")
prs.save(OUT + "03-klaravex-2.0-growth-os.pptx")

# ------------------------------------------------------ Deck 4: KLARAVEX-OS
prs = new_deck()
title_slide(prs, "INTERNAL — KLARAVEX-OS",
            "One screen for the whole company",
            "The internal operating console: delivery, pipeline, finances, comms, the "
            "knowledge layer, and the AI agent fleet — real Postgres, real connectors, "
            "honest status. Single-operator by design.", dark=True)
bullets_slide(prs, "What lives on the screen", [
    ("Delivery", "Client work and project state — what is due, what is blocked, what shipped."),
    ("Pipeline & finances", "Prospects, proposals, MRR, and billing state in one view."),
    ("Comms & knowledge layer", "Conversations and the durable memory the agents read and write."),
    ("Agent fleet", "Every autonomous stream with honest status — running, idle, failed — no green-washing."),
])
flow_slide(prs, "How the systems tie together", [
    ("klaravex (production)", "The live monolith: support runtime, voice/SMS/chat, billing, website, autonomous engineering loop."),
    ("Klaravex2.0 (Growth OS)", "Layers A + C: revenue-agent charters, outbox, Growth API :4210, systemd cadence."),
    ("KLARAVEX-OS :4100", "Layer D. Reads the operating picture; drives growth exclusively through the Growth API — never around it."),
    ("Operator", "One owner approves gated publishes, watches scorecards, and intervenes where judgment is needed."),
], footer="Rule of the architecture: the cockpit talks to control planes, not to internals.")
bullets_slide(prs, "Design decisions that matter", [
    ("Single-operator, single-tenant", "One instance, one owner. The multi-tenant product build is a separate, later phase."),
    ("Honest status over dashboards-theater", "Real Postgres, real connectors; if a connector is down, the screen says so."),
    ("Separate from the client portal", "Founders OS — the client-facing dashboard (readiness scores, compliance tracker, magic-link login) — is a different product that shares the philosophy, not the instance."),
    ("Credentials via vault", "Runs with 1Password-injected environment — no secrets in files."),
])
bullets_slide(prs, "The full estate at a glance", [
    ("klaravex.com", "B2B managed IT & security — AI-first delivery, senior judgment on top (deck 1)."),
    ("personal.klaravex.com", "Consumer help by Klaravex AI — $29 sessions, free scam recovery (deck 2)."),
    ("Klaravex 2.0", "Growth OS strangler-fig — autonomous revenue streams on independent cadence (deck 3)."),
    ("KLARAVEX-OS", "The cockpit that watches and steers all of it — this deck."),
])
cta_slide(prs, "One company. Four surfaces. One screen.",
          ["Production: real connectors, honest status, gated automation.",
           "Next: wire Layer D to the Growth API for every stream."],
          "KLARAVEX-OS — internal · localhost:4100")
prs.save(OUT + "04-klaravex-os.pptx")

print("done")
for f in ["01-klaravex-b2b", "02-klaravex-personal",
          "03-klaravex-2.0-growth-os", "04-klaravex-os"]:
    p = Presentation(OUT + f + ".pptx")
    print(f, "slides:", len(p.slides.__iter__.__self__._sldIdLst))
